"""
Document Processor

Handles:
- HTML text extraction and cleaning
- PDF text extraction
- PowerPoint text extraction
- Text chunking with metadata preservation
"""

import io
import re
import logging
from typing import List, Dict
from bs4 import BeautifulSoup
from src.config import CHUNK_SIZE, CHUNK_OVERLAP, PDF_EXTRACTION_TIMEOUT

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Processes documents from various sources into chunks with metadata."""

    def __init__(self, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP):
        """
        Initialize document processor.

        Args:
            chunk_size: Number of words per chunk
            chunk_overlap: Number of overlapping words between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def extract_html_text(self, html_content: str) -> str:
        """
        Extract clean text from HTML content.

        Args:
            html_content: HTML string

        Returns:
            Clean text content
        """
        if not html_content:
            return ""

        try:
            # Try lxml parser first (more forgiving)
            soup = BeautifulSoup(html_content, 'lxml')
        except Exception as e:
            logger.warning(f"lxml parsing failed, trying html.parser: {e}")
            try:
                soup = BeautifulSoup(html_content, 'html.parser')
            except Exception as e2:
                logger.error(f"HTML parsing completely failed: {e2}")
                # Fallback: strip HTML tags manually
                return re.sub('<[^<]+?>', '', html_content)

        # Remove script and style elements
        for script in soup(['script', 'style', 'meta', 'link']):
            script.decompose()

        # Special handling for tables to preserve structure
        for table in soup.find_all('table'):
            # Convert table to readable text format
            rows = []
            for row in table.find_all('tr'):
                cells = [cell.get_text(strip=True) for cell in row.find_all(['td', 'th'])]
                if cells:  # Only add non-empty rows
                    rows.append(' | '.join(cells))

            if rows:
                # Replace table with formatted text
                table_text = '\n'.join(rows)
                table.replace_with(soup.new_string(f"\n{table_text}\n"))

        # Get text and clean whitespace
        text = soup.get_text(separator='\n')

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return text

    def extract_pdf_text(self, pdf_bytes: bytes) -> str:
        """
        Extract text from PDF file.

        Args:
            pdf_bytes: PDF file content as bytes

        Returns:
            Extracted text
        """
        try:
            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
            text = ""

            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                except Exception as e:
                    logger.warning(f"Error extracting PDF page {page_num}: {e}")
                    continue

            if not text.strip():
                raise ValueError("No text extracted from PDF")

            return text

        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ""

    def extract_pptx_text(self, pptx_bytes: bytes) -> str:
        """
        Extract text from PowerPoint file.

        Args:
            pptx_bytes: PPTX file content as bytes

        Returns:
            Extracted text
        """
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(pptx_bytes))
            text = ""

            for slide_num, slide in enumerate(prs.slides):
                try:
                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                except Exception as e:
                    logger.warning(f"Error extracting PPTX slide {slide_num}: {e}")
                    continue

            if not text.strip():
                raise ValueError("No text extracted from PowerPoint")

            return text

        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            return ""

    def chunk_text(self, text: str, metadata: Dict) -> List[Dict]:
        """
        Chunk text into overlapping segments.

        Args:
            text: Text to chunk
            metadata: Metadata to attach to each chunk

        Returns:
            List of chunk dicts with text and metadata
        """
        if not text or not text.strip():
            return []

        # Split into words
        words = text.split()

        if len(words) == 0:
            return []

        # If text is shorter than chunk size, return as single chunk
        if len(words) <= self.chunk_size:
            return [{
                'text': text,
                'metadata': {
                    **metadata,
                    'chunk_index': 0,
                    'total_chunks': 1
                }
            }]

        chunks = []
        start = 0

        while start < len(words):
            # Get chunk of words
            end = min(start + self.chunk_size, len(words))
            chunk_words = words[start:end]
            chunk_text = ' '.join(chunk_words)

            chunks.append({
                'text': chunk_text,
                'metadata': {
                    **metadata,
                    'chunk_index': len(chunks),
                    'total_chunks': -1  # Will be updated after loop
                }
            })

            # Move to next chunk with overlap
            if end >= len(words):
                break
            start += self.chunk_size - self.chunk_overlap

        # Update total_chunks in all metadata
        for chunk in chunks:
            chunk['metadata']['total_chunks'] = len(chunks)

        return chunks

    def process_html_content(self, html_content: str, metadata: Dict) -> List[Dict]:
        """
        Process HTML content into chunks.

        Args:
            html_content: HTML string
            metadata: Metadata to attach

        Returns:
            List of chunks
        """
        text = self.extract_html_text(html_content)
        return self.chunk_text(text, metadata)

    def process_pdf_content(self, pdf_bytes: bytes, metadata: Dict) -> List[Dict]:
        """
        Process PDF content into chunks.

        Args:
            pdf_bytes: PDF file bytes
            metadata: Metadata to attach

        Returns:
            List of chunks
        """
        text = self.extract_pdf_text(pdf_bytes)
        if not text:
            logger.warning(f"No text extracted from PDF: {metadata.get('title', 'Unknown')}")
            return []
        return self.chunk_text(text, metadata)

    def process_pptx_content(self, pptx_bytes: bytes, metadata: Dict) -> List[Dict]:
        """
        Process PowerPoint content into chunks.

        Args:
            pptx_bytes: PPTX file bytes
            metadata: Metadata to attach

        Returns:
            List of chunks
        """
        text = self.extract_pptx_text(pptx_bytes)
        if not text:
            logger.warning(f"No text extracted from PowerPoint: {metadata.get('title', 'Unknown')}")
            return []
        return self.chunk_text(text, metadata)

    def process_document(self, content: Dict, content_type: str) -> List[Dict]:
        """
        Process a document based on its type.

        Args:
            content: Content dict with 'data' and optional 'bytes'/'html' keys
            content_type: Type of content (file, page, assignment, etc.)

        Returns:
            List of chunks
        """
        if content_type == 'file':
            # File content - check mime type
            mime_type = content.get('mime_type', content.get('content-type', ''))
            file_bytes = content.get('bytes', b'')

            if 'pdf' in mime_type.lower():
                return self.process_pdf_content(file_bytes, content.get('metadata', {}))
            elif 'powerpoint' in mime_type.lower() or 'presentation' in mime_type.lower():
                return self.process_pptx_content(file_bytes, content.get('metadata', {}))
            else:
                logger.warning(f"Unsupported file type: {mime_type}")
                return []
        else:
            # HTML content (page, assignment, discussion, etc.)
            html_content = content.get('html', content.get('body', ''))
            return self.process_html_content(html_content, content.get('metadata', {}))
